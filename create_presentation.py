"""
Генератор презентации HSE AI Analytics — KMG Hackathon 2026.
Запуск: python create_presentation.py
Результат: HSE_Presentation.pptx
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    import subprocess, sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ── Цветовая палитра ──────────────────────────────────────────────────────────
C_DARK_BLUE  = RGBColor(0x1e, 0x3a, 0x8a)
C_BLUE       = RGBColor(0x2d, 0x6a, 0xd4)
C_RED        = RGBColor(0xef, 0x44, 0x44)
C_AMBER      = RGBColor(0xf5, 0x9e, 0x0b)
C_GREEN      = RGBColor(0x22, 0xc5, 0x5e)
C_WHITE      = RGBColor(0xff, 0xff, 0xff)
C_LIGHT      = RGBColor(0xf1, 0xf5, 0xf9)
C_GRAY       = RGBColor(0x64, 0x74, 0x8b)
C_DARK       = RGBColor(0x1e, 0x29, 0x3b)

W = Inches(13.33)   # широкоформатный слайд 16:9
H = Inches(7.5)


# ── Хелперы ───────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # полностью пустой
    return prs.slides.add_slide(layout)


def rect(slide, x, y, w, h, fill_rgb: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def txt(slide, text: str, x, y, w, h,
        font_size=18, bold=False, color: RGBColor = C_WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bullet(tf, text: str, level=0, size=16,
               color: RGBColor = C_DARK, bold=False):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold


def header_bar(slide, title: str, subtitle: str = ""):
    """Синяя полоса сверху с заголовком."""
    rect(slide, 0, 0, W, Inches(1.3), C_DARK_BLUE)
    txt(slide, title, Inches(0.4), Inches(0.15), Inches(10), Inches(0.7),
        font_size=28, bold=True, color=C_WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.4), Inches(0.82), Inches(10), Inches(0.4),
            font_size=14, color=RGBColor(0xba, 0xd4, 0xf8))


def footer(slide, page_num: int, total: int):
    """Нижняя строка со ссылкой на Streamlit + номер слайда."""
    txt(slide, "http://89.207.255.254:8502",
        Inches(0.3), Inches(7.1), Inches(5), Inches(0.3),
        font_size=10, color=C_GRAY, italic=True)
    txt(slide, f"{page_num} / {total}",
        Inches(12.0), Inches(7.1), Inches(1.2), Inches(0.3),
        font_size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)


def kpi_box(slide, value: str, label: str,
            x, y, w=Inches(2.5), h=Inches(1.4),
            bg: RGBColor = C_DARK_BLUE):
    rect(slide, x, y, w, h, bg)
    txt(slide, value, x, y + Inches(0.1), w, Inches(0.7),
        font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + Inches(0.8), w, Inches(0.5),
        font_size=12, color=RGBColor(0xba, 0xd4, 0xf8), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙДЫ
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Слайд 1 — Титульный."""
    sl = blank_slide(prs)

    # Фон — тёмно-синий градиент (два прямоугольника)
    rect(sl, 0, 0, W, H, C_DARK_BLUE)
    rect(sl, 0, Inches(5.5), W, Inches(2.0), C_DARK)

    # Декоративная полоса
    rect(sl, 0, Inches(5.3), W, Inches(0.06), C_RED)

    # Логотип-заглушка (щит)
    txt(sl, "🛡️", Inches(0.5), Inches(0.6), Inches(1.5), Inches(1.5),
        font_size=64, color=C_WHITE)

    # Заголовок
    txt(sl, "HSE AI Analytics",
        Inches(2.0), Inches(0.5), Inches(10), Inches(1.3),
        font_size=52, bold=True, color=C_WHITE)

    txt(sl, "AI-аналитика охраны труда и промышленной безопасности",
        Inches(2.0), Inches(1.75), Inches(10), Inches(0.7),
        font_size=22, color=RGBColor(0xba, 0xd4, 0xf8))

    # Разделитель
    rect(sl, Inches(2.0), Inches(2.5), Inches(4.0), Inches(0.05), C_AMBER)

    txt(sl, "KMG Hackathon 2026",
        Inches(2.0), Inches(2.7), Inches(6), Inches(0.5),
        font_size=18, bold=True, color=C_AMBER)

    # Нижняя плашка
    txt(sl, "Предиктивный инструмент управления безопасностью с измеримым экономическим эффектом",
        Inches(0.5), Inches(5.8), Inches(12), Inches(0.6),
        font_size=14, color=RGBColor(0x94, 0xa3, 0xb8), align=PP_ALIGN.CENTER)

    txt(sl, "FastAPI  •  Prophet  •  Gemini 2.0 Flash  •  Streamlit  •  Docker",
        Inches(0.5), Inches(6.4), Inches(12), Inches(0.5),
        font_size=13, color=RGBColor(0x64, 0x74, 0x8b), align=PP_ALIGN.CENTER)
    footer(sl, 1, TOTAL)


def slide_problem(prs):
    """Слайд 2 — Проблема."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "Проблема: реактивный подход к безопасности",
               "Текущее состояние охраны труда в нефтегазовой отрасли")

    # 3 KPI-бокса с правильными отступами
    cols = [
        ("18 НС/год",      "Несчастных случаев",        C_RED),
        ("120 микротравм",  "Ежегодно",                  C_AMBER),
        ("72 часа",         "Среднее время реагирования", C_BLUE),
    ]
    for i, (val, lbl, color) in enumerate(cols):
        x = Inches(0.3 + i * 4.25)
        kpi_box(sl, val, lbl, x, Inches(1.5), Inches(3.95), Inches(1.4), bg=color)

    # Ключевые проблемы
    rect(sl, Inches(0.3), Inches(3.1), Inches(12.73), Inches(3.4), C_WHITE)

    txt(sl, "Ключевые проблемы", Inches(0.5), Inches(3.2), Inches(6), Inches(0.4),
        font_size=16, bold=True, color=C_DARK_BLUE)

    items = [
        "Нет предиктивной аналитики — инциденты расследуются постфактум",
        "Данные Карт Коргау и происшествий не связаны между собой",
        "Ручной анализ сотен наблюдений занимает дни",
        "Экономический эффект мер безопасности не рассчитывается",
        "350 опасных ситуаций в год остаются без превентивных мер",
    ]
    for i, item in enumerate(items):
        ty = Inches(3.65 + i * 0.55)
        txt(sl, f"❌  {item}", Inches(0.5), ty, Inches(12), Inches(0.45),
            font_size=13, color=C_DARK)

    # Итоговый баннер
    rect(sl, Inches(0.3), Inches(6.65), Inches(12.73), Inches(0.5), C_RED)
    txt(sl, "Потери: ~200 млн ₸/год   |   Цена бездействия растёт",
        Inches(0.3), Inches(6.68), Inches(12.73), Inches(0.45),
        font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    footer(sl, 2, TOTAL)


def slide_solution(prs):
    """Слайд 3 — Решение."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "Решение: HSE AI Analytics",
               "AI-аналитический слой поверх двух модулей HSE-системы")

    # Два модуля
    for i, (title, items, color) in enumerate([
        ("Компонент A — Инциденты", [
            "NLP-классификация описаний (Gemini)",
            "Паттерны по типу, орг., локации, смене",
            "Кластеризация причин (TF-IDF + KMeans)",
            "Топ-5 рекомендаций от AI",
            "Prophet-прогноз + сценарный анализ",
            "CV-анализ фото (Gemini Vision)",
        ], C_DARK_BLUE),
        ("Компонент B — Карты Коргау", [
            "Паттерны нарушений по орг./категориям",
            "4-уровневая система алертов",
            "Рейтинг организаций по compliance",
            "Корреляция нарушений с инцидентами",
            "Pre-alert при росте нарушений",
            "AI-кластеризация наблюдений",
        ], C_BLUE),
    ]):
        x = Inches(0.3 + i * 6.6)
        rect(sl, x, Inches(1.45), Inches(6.3), Inches(4.8), color)
        txt(sl, title, x + Inches(0.15), Inches(1.5), Inches(6.0), Inches(0.6),
            font_size=16, bold=True, color=C_WHITE)
        rect(sl, x, Inches(2.1), Inches(6.3), Inches(0.03), C_WHITE)

        txb = sl.shapes.add_textbox(x + Inches(0.15), Inches(2.2),
                                    Inches(6.0), Inches(3.8))
        tf = txb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.add_run().text = ""  # пустой первый параграф
        for item in items:
            add_bullet(tf, f"✓  {item}", size=14, color=C_WHITE)

    # Стрелка объединения
    txt(sl, "+", Inches(6.15), Inches(3.5), Inches(1.0), Inches(0.8),
        font_size=36, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

    # Итог внизу
    rect(sl, Inches(0.3), Inches(6.35), Inches(12.7), Inches(0.9), C_DARK_BLUE)
    txt(sl, "15 REST-эндпоинтов  •  Streamlit-дашборд  •  PDF-отчёт  •  Сценарный анализ  •  Docker",
        Inches(0.3), Inches(6.35), Inches(12.7), Inches(0.9),
        font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    footer(sl, 3, TOTAL)


def slide_architecture(prs):
    """Слайд 4 — Архитектура."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "Трёхуровневая архитектура",
               "ETL → ML/AI Core → Представление")

    # Уровни
    levels = [
        ("📂  Уровень данных",
         "data_loader.py  —  ETL Pipeline\nincidents.xlsx (220 строк, 50 колонок)\nkorgau_cards.xlsx (9 917 строк, 16 колонок)\nОбогащение: смена, стаж, флаги, resolved",
         C_DARK_BLUE),
        ("🧠  AI / ML Core + NLP",
         "Prophet  —  time-series forecast\nTF-IDF + KMeans  —  кластеры причин\nRisk Scoring  —  формула из ТЗ\nPearson / Spearman  —  корреляция\nGemini 2.0 Flash  —  NLP + Vision",
         C_BLUE),
        ("🖥️  Представление",
         "Streamlit  —  4 таба, 10 Plotly-графиков\nFastAPI  —  15 эндпоинтов, OpenAPI /docs\nfpdf2  —  PDF-отчёт 4 страницы\nAlert Manager  —  4 уровня CRITICAL→LOW",
         C_RED),
    ]

    for i, (title, body, color) in enumerate(levels):
        y = Inches(1.5 + i * 1.9)
        rect(sl, Inches(0.4), y, Inches(12.5), Inches(1.75), color)
        txt(sl, title, Inches(0.6), y + Inches(0.1), Inches(4.5), Inches(0.55),
            font_size=16, bold=True, color=C_WHITE)
        txt(sl, body, Inches(5.2), y + Inches(0.08), Inches(7.5), Inches(1.5),
            font_size=13, color=C_WHITE)

        # Разделитель
        if i < 2:
            txt(sl, "▼", Inches(6.3), y + Inches(1.78), Inches(0.8), Inches(0.3),
                font_size=14, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)


def slide_predict(prs):
    """Слайд 5 — Предиктивная аналитика."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "Предиктивная аналитика",
               "Prophet Forecast  •  Risk Scoring  •  Сценарный анализ мер контроля")

    # Колонка 1 — Prophet
    rect(sl, Inches(0.3), Inches(1.4), Inches(4.1), Inches(5.8), C_WHITE)
    txt(sl, "🔮  Prophet Forecast", Inches(0.4), Inches(1.5), Inches(4.0), Inches(0.5),
        font_size=15, bold=True, color=C_DARK_BLUE)
    txb = sl.shapes.add_textbox(Inches(0.4), Inches(2.05), Inches(3.9), Inches(4.8))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.add_run().text = ""
    for item in [
        "Прогноз на 12 / 6 / 3 месяца",
        "95% доверительный интервал",
        "Годовая сезонность",
        "Fallback: линейная регрессия",
        "",
        "Метрики backtesting:",
        "  MAE — абсолютная ошибка",
        "  RMSE — среднеквадратичная",
        "  MAPE — процентная ошибка",
    ]:
        add_bullet(tf, f"• {item}" if item and not item.startswith(" ") else item,
                   size=13 if not item.startswith("  ") else 11,
                   color=C_DARK if not item.startswith("Метрики") else C_GREEN,
                   bold=item.startswith("Метрики"))

    # Колонка 2 — Risk Scoring
    rect(sl, Inches(4.65), Inches(1.4), Inches(4.1), Inches(5.8), C_WHITE)
    txt(sl, "⚡  Risk Scoring (ТЗ)", Inches(4.75), Inches(1.5), Inches(4.0), Inches(0.5),
        font_size=15, bold=True, color=C_DARK_BLUE)

    formula_lines = [
        "risk_score =",
        "  0.35 × incident_rate",
        "+ 0.30 × violation_rate",
        "+ 0.20 × trend_slope",
        "+ 0.15 × severity_factor",
        "",
        "Нормализация: 0–100",
        "CRITICAL ≥ 70",
        "HIGH     ≥ 50",
        "MEDIUM   ≥ 25",
        "LOW       < 25",
    ]
    txb2 = sl.shapes.add_textbox(Inches(4.75), Inches(2.05), Inches(3.9), Inches(4.8))
    tf2 = txb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.add_run().text = ""
    for line in formula_lines:
        add_bullet(tf2, line, size=12,
                   color=C_RED if line.startswith("CRITICAL") else
                         C_AMBER if line.startswith("HIGH") else
                         C_DARK)

    # Колонка 3 — Scenario
    rect(sl, Inches(9.0), Inches(1.4), Inches(4.1), Inches(5.8), C_WHITE)
    txt(sl, "🎯  Сценарный анализ", Inches(9.1), Inches(1.5), Inches(4.0), Inches(0.5),
        font_size=15, bold=True, color=C_DARK_BLUE)
    txb3 = sl.shapes.add_textbox(Inches(9.1), Inches(2.05), Inches(3.9), Inches(4.8))
    tf3 = txb3.text_frame; tf3.word_wrap = True
    p3 = tf3.paragraphs[0]; p3.add_run().text = ""
    for item in [
        "9 мер контроля (LOTO,",
        "  СИЗ, Высота, Огонь...)",
        "× Pearson r = 0.415",
        "= ожид. снижение НС",
        "Экономика: × 5 млн ₸",
        "",
        "Кластеры причин:",
        "TF-IDF + KMeans (n=6)",
        "Silhouette = 0.835",
        "Davies-Bouldin = 1.204",
    ]:
        add_bullet(tf3, item, size=13 if not item.startswith("Silhouette") else 12,
                   color=C_GREEN if "0.835" in item or "Кластеры" in item else C_DARK,
                   bold="Кластеры" in item or "Silhouette" in item)


def slide_korgau(prs):
    """Слайд 6 — Карты Коргау."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "Карты Коргау — поведенческий аудит безопасности",
               "9 917 наблюдений  •  Корреляция с инцидентами r = 0.415 (Pearson)")

    features = [
        ("F-10", "Паттерны нарушений", "groupby (org, category, month)\nтоп нарушителей и категорий"),
        ("F-11", "Система алертов", "CRITICAL > 50 нарушений/30 дней\nHIGH > 30  •  MEDIUM > 15  •  LOW > 5"),
        ("F-12", "Рейтинг организаций", "64 org  •  violation_rate %\ncompliance_rate  •  work_stops"),
        ("F-13", "Корреляция НС", "Pearson r=0.415  Spearman r=0.501\nстатистически значима"),
        ("F-14", "Pre-alert тренд", "рост / снижение / стабильно\nпо последним 3 месяцам"),
        ("F-15", "AI-кластеризация", "Gemini 2.0 Flash\n10 тематических кластеров"),
    ]

    for i, (fid, title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        x = Inches(0.3 + col * 4.35)
        y = Inches(1.5 + row * 2.55)
        rect(sl, x, y, Inches(4.1), Inches(2.3), C_DARK_BLUE if row == 0 else C_BLUE)
        txt(sl, fid, x + Inches(0.15), y + Inches(0.1), Inches(1.0), Inches(0.35),
            font_size=11, bold=True, color=C_AMBER)
        txt(sl, title, x + Inches(0.15), y + Inches(0.4), Inches(3.8), Inches(0.45),
            font_size=14, bold=True, color=C_WHITE)
        txt(sl, desc, x + Inches(0.15), y + Inches(0.85), Inches(3.8), Inches(1.3),
            font_size=12, color=RGBColor(0xba, 0xd4, 0xf8))


def slide_economic(prs):
    """Слайд 7 — Экономический эффект."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "Экономический эффект",
               "Измеримый результат внедрения AI-аналитики")

    # KPI в людях
    rect(sl, Inches(0.3), Inches(1.4), Inches(12.7), Inches(0.4), C_DARK_BLUE)
    txt(sl, "В ЛЮДЯХ", Inches(0.3), Inches(1.4), Inches(12.7), Inches(0.4),
        font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    kpis = [
        ("↓ 38%",  "Несчастных\nслучаев"),
        ("↓ 40%",  "Микро-\nтравм"),
        ("↓ 50%",  "Опасных\nситуаций"),
        ("↓ 83%",  "Время\nреагирования"),
        ("↑ 70%",  "Охват\nпредиктивом"),
    ]
    for i, (val, lbl) in enumerate(kpis):
        x = Inches(0.3 + i * 2.56)
        color = C_GREEN if val.startswith("↑") else C_RED
        kpi_box(sl, val, lbl, x, Inches(1.85), Inches(2.45), Inches(1.35), bg=color)

    # В деньгах — таблица
    rect(sl, Inches(0.3), Inches(3.35), Inches(12.7), Inches(0.4), C_DARK_BLUE)
    txt(sl, "В ДЕНЬГАХ", Inches(0.3), Inches(3.35), Inches(12.7), Inches(0.4),
        font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("Прямые затраты на НС  (7 × 5 млн ₸)",          "≈ 35 000 000 ₸"),
        ("Косвенные потери  (простой, репутация)",         "≈ 70 000 000 ₸"),
        ("Штрафы и предписания  (−30%)",                   "≈  5 000 000 ₸"),
        ("Расследования и отчётность  (−70%)",             "≈  8 000 000 ₸"),
        ("ROI автоматизации аудитов  (+40%)",              "≈  3 000 000 ₸"),
    ]
    for i, (label, value) in enumerate(rows):
        y = Inches(3.8 + i * 0.52)
        bg = C_WHITE if i % 2 == 0 else C_LIGHT
        rect(sl, Inches(0.3), y, Inches(12.7), Inches(0.5), bg)
        txt(sl, label, Inches(0.5), y + Inches(0.08), Inches(9.5), Inches(0.35),
            font_size=13, color=C_DARK)
        txt(sl, value, Inches(10.0), y + Inches(0.08), Inches(2.8), Inches(0.35),
            font_size=13, bold=True, color=C_DARK_BLUE, align=PP_ALIGN.RIGHT)

    # Итого
    rect(sl, Inches(0.3), Inches(6.45), Inches(12.7), Inches(0.7), C_GREEN)
    txt(sl, "ИТОГО:   ≈ 121 000 000 ₸  (~250 000 USD)  в год",
        Inches(0.3), Inches(6.45), Inches(12.7), Inches(0.7),
        font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    footer(sl, 7, TOTAL)


def slide_tech(prs):
    """Слайд 8 — Технологический стек."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "Технологический стек",
               "Production-ready Python-монолит с Docker Compose")

    stack = [
        ("🔧  Backend",    "Python 3.11  •  FastAPI  •  uvicorn\n15 REST-эндпоинтов  •  OpenAPI /docs",    C_DARK_BLUE),
        ("📊  ML Core",    "Prophet (time series)\nscikit-learn: TF-IDF + KMeans\nRisk Scoring + Pearson/Spearman",  C_BLUE),
        ("🤖  NLP / CV",   "Google Gemini 2.0 Flash\nклассификация  •  рекомендации\nGemini Vision (CV-анализ фото)",C_BLUE),
        ("📈  Frontend",   "Streamlit  •  4 таба\n10 Plotly-графиков\nCI-лента прогноза",                  C_DARK_BLUE),
        ("📄  Отчёты",     "fpdf2  •  4-страничный PDF\nDejaVu TTF (кириллица)\nKPI + графики + алерты",   C_RED),
        ("🐳  DevOps",     "Docker + Docker Compose\nhse-api :8002  •  hse-ui :8502\nhealthcheck + networks",C_RED),
    ]

    for i, (title, body, color) in enumerate(stack):
        col = i % 3
        row = i // 3
        x = Inches(0.3 + col * 4.35)
        y = Inches(1.5 + row * 2.85)
        rect(sl, x, y, Inches(4.1), Inches(2.6), color)
        txt(sl, title, x + Inches(0.15), y + Inches(0.12), Inches(3.8), Inches(0.45),
            font_size=15, bold=True, color=C_WHITE)
        rect(sl, x, y + Inches(0.58), Inches(4.1), Inches(0.03), C_WHITE)
        txt(sl, body, x + Inches(0.15), y + Inches(0.68), Inches(3.8), Inches(1.6),
            font_size=13, color=RGBColor(0xba, 0xd4, 0xf8))


def slide_demo(prs):
    """Слайд 9 — Дашборд / скриншоты."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "Дашборд — 4 вкладки Streamlit",
               "http://localhost:8502  •  Swagger http://localhost:8002/docs")

    tabs = [
        ("📊 Инциденты",
         ["Динамика + скользящее среднее 3м",
          "Pie по типам инцидентов",
          "Топ-10 организаций",
          "Кластеры причин (TF-IDF + KMeans)",
          "NLP-классификатор (Gemini)",
          "CV-анализ фото (Gemini Vision)"],
         C_DARK_BLUE),
        ("🔍 Карты Коргау",
         ["Тренд нарушений топ-5 org",
          "Рейтинг организаций (% нарушений)",
          "Хорошие vs плохие практики",
          "Pearson r / Spearman r метрики"],
         C_BLUE),
        ("🔮 Предикт",
         ["Prophet: прогноз + 95% CI лента",
          "Риск-скоринг + heatmap (0–100)",
          "Матрица корреляций (heatmap)",
          "Сценарный анализ мер контроля"],
         C_BLUE),
        ("🚨 Алерты",
         ["CRITICAL/HIGH/MEDIUM/LOW карточки",
          "Фильтры по org и уровню",
          "Счётчики по уровням",
          "Кнопка скачать PDF-отчёт"],
         C_RED),
    ]

    for i, (title, items, color) in enumerate(tabs):
        x = Inches(0.3 + i * 3.2)
        rect(sl, x, Inches(1.4), Inches(3.0), Inches(5.8), color)
        txt(sl, title, x + Inches(0.1), Inches(1.5), Inches(2.8), Inches(0.5),
            font_size=14, bold=True, color=C_WHITE)
        rect(sl, x, Inches(2.0), Inches(3.0), Inches(0.03), C_WHITE)
        txb = sl.shapes.add_textbox(x + Inches(0.1), Inches(2.1), Inches(2.8), Inches(4.8))
        tf = txb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.add_run().text = ""
        for item in items:
            add_bullet(tf, f"✓  {item}", size=12, color=C_WHITE)


def slide_results(prs):
    """Слайд 10 — Итог / CTA."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_DARK_BLUE)
    rect(sl, 0, Inches(5.6), W, Inches(1.9), C_DARK)
    rect(sl, 0, Inches(5.5), W, Inches(0.06), C_RED)

    txt(sl, "🛡️", Inches(0.5), Inches(0.5), Inches(1.2), Inches(1.2),
        font_size=56, color=C_WHITE)

    txt(sl, "Результат",
        Inches(1.8), Inches(0.5), Inches(10), Inches(0.9),
        font_size=46, bold=True, color=C_WHITE)

    checklist = [
        "✅  Рабочий AI-дашборд  —  аналитика прошлых инцидентов",
        "✅  Предикт на 12 месяцев  —  Prophet + 95% CI",
        "✅  Сценарный анализ  —  расчёт эффекта мер контроля",
        "✅  4-уровневые алерты  —  CRITICAL → LOW",
        "✅  Топ-5 AI-рекомендаций  —  Gemini 2.0 Flash",
        "✅  CV-анализ фото  —  Gemini Vision",
        "✅  PDF-отчёт  —  4 страницы, кириллица",
        "✅  Экономика  —  121 млн ₸ / год",
    ]
    txb = sl.shapes.add_textbox(Inches(1.8), Inches(1.5), Inches(11), Inches(3.8))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.add_run().text = ""
    for item in checklist:
        add_bullet(tf, item, size=15,
                   color=C_WHITE if "✅" in item else RGBColor(0xba, 0xd4, 0xf8))

    txt(sl, "FastAPI :8002  •  Streamlit :8502  •  docker compose up --build",
        Inches(0.5), Inches(5.8), Inches(12), Inches(0.5),
        font_size=14, color=RGBColor(0xba, 0xd4, 0xf8), align=PP_ALIGN.CENTER)

    txt(sl, "KMG Hackathon 2026  •  github.com/Islambek201632838/kmg-hse",
        Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
        font_size=13, color=RGBColor(0x64, 0x74, 0x8b), align=PP_ALIGN.CENTER)
    footer(sl, 10, TOTAL)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def slide_criteria(prs):
    """Обоснование по 6 критериям хакатона."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "Критерии хакатона — как мы их закрываем", "Раздел 6 ТЗ: 6 критериев оценки")

    criteria = [
        (C_DARK_BLUE, "25%", "Точность предиктива",
         "Prophet + 95% CI · backtesting на истории\nПрогноз 12/6/3 мес · годовая сезонность"),
        (C_BLUE, "20%", "Качество рекомендаций",
         "Gemini топ-5 с приоритетом и обоснованием\nНа основе реальных паттернов (орг, тип, причины)"),
        (C_GREEN, "15%", "UX дашборда",
         "Streamlit 4 таба · 10+ Plotly-графиков\nParallel loading · KPI-карточки · фильтры"),
        (C_AMBER, "15%", "Система алертов",
         "4 уровня: CRITICAL (x2) · HIGH (>3 раз/тип)\nMEDIUM (YoY >15%) · LOW (информация)"),
        (RGBColor(0x6a, 0x1b, 0x9a), "15%", "Интеграция API",
         "FastAPI · 17 REST эндпоинтов · OpenAPI /docs\nDocker Compose · JSON · CORS"),
        (C_RED, "10%", "Экономический эффект",
         "Сценарное моделирование: 9 мер контроля\nPearson r=0.415 · ~121 млн ₸/год экономии"),
    ]

    for i, (color, weight, title, body) in enumerate(criteria):
        col, row = i % 3, i // 3
        lx = Inches(0.3 + col * 4.3)
        ty = Inches(1.5 + row * 2.85)

        rect(sl, lx, ty, Inches(4.0), Inches(2.6), C_WHITE)
        rect(sl, lx, ty, Inches(4.0), Inches(0.5), color)

        txt(sl, weight, lx + Inches(0.1), ty + Inches(0.05), Inches(0.8), Inches(0.4),
            font_size=20, bold=True, color=C_WHITE)
        txt(sl, title, lx + Inches(1.0), ty + Inches(0.08), Inches(2.8), Inches(0.35),
            font_size=13, bold=True, color=C_WHITE)
        txt(sl, body, lx + Inches(0.15), ty + Inches(0.6), Inches(3.7), Inches(1.8),
            font_size=11, color=C_DARK)
    footer(sl, 8, TOTAL)


def slide_optimizations(prs):
    """Оптимизации и улучшения."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "Оптимизации", "Параллельность · кеш · алерт-правила · экспорт")

    items = [
        (C_GREEN, "ThreadPoolExecutor",
         "Параллельная загрузка API-запросов в каждом табе\nTab 2: 4 запроса параллельно (trends + rankings + gvb + correlation)"),
        (C_BLUE, "Кеш 10 мин + LRU",
         "Streamlit @cache_data TTL=600s для API\n@lru_cache для Excel-парсинга (грузится один раз)"),
        (C_AMBER, "Алерт: тип >3x за 30 дней",
         "HIGH-алерт если один тип нарушения повторяется >3 раз\nАвтоматическое повышение уровня + причина в ответе"),
        (C_RED, "Алерт: YoY рост >15%",
         "MEDIUM-алерт если нарушения выросли на >15% к прошлому году\nСравнение 365 дней текущий vs предыдущий период"),
        (RGBColor(0x6a, 0x1b, 0x9a), "Excel-экспорт (F-09)",
         "5 листов: KPI, Alerts, Rankings, Risk Scores, Forecast\nopenpyxl + auto-width + color headers"),
        (C_DARK_BLUE, "Korgau classify API (F-15)",
         "POST /api/korgau/classify — AI-классификация наблюдений\n10 кластеров: СИЗ, высота, LOTO, пожарная и др."),
    ]

    for i, (color, title, body) in enumerate(items):
        col, row = i % 2, i // 2
        lx = Inches(0.3 + col * 6.52)
        ty = Inches(1.5 + row * 1.85)

        rect(sl, lx, ty, Inches(6.22), Inches(1.65), C_WHITE)
        rect(sl, lx, ty, Inches(0.1), Inches(1.65), color)
        txt(sl, title, lx + Inches(0.2), ty + Inches(0.1), Inches(5.8), Inches(0.35),
            font_size=13, bold=True, color=color)
        txt(sl, body, lx + Inches(0.2), ty + Inches(0.5), Inches(5.8), Inches(1.0),
            font_size=11, color=C_DARK)


def slide_ml_pipeline(prs):
    """ML/AI пайплайн — что, зачем, как работает."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "ML/AI пайплайн — как система предсказывает",
               "От сырых данных к предиктивным решениям")

    # Pipeline flow: 5 steps
    steps = [
        (C_DARK_BLUE, "1. ETL",
         "xlsx → pandas\nОчистка: NaT, пропуски\nОбогащение: смены,\nстаж, тип, тяжесть"),
        (C_BLUE, "2. Feature\nEngineering",
         "Агрегация по месяцам\nRolling avg (3 мес)\nMoM % изменение\nСмена: день/вечер/ночь"),
        (RGBColor(0x6a, 0x1b, 0x9a), "3. ML Models",
         "Prophet: тренд + сезонность\nKMeans: кластеры причин\nPearson/Spearman: связь\nRisk Scoring: формула"),
        (C_RED, "4. NLP / CV",
         "Gemini Flash:\n• классификация инцидентов\n• рекомендации (топ-5)\n• анализ фото (Vision)"),
        (C_GREEN, "5. Predict",
         "Прогноз на 12 мес + 95% CI\nРиск-скор 0–100 по орг\nСценарий: мера → экономия\nАлерты: 4 уровня"),
    ]

    for i, (color, title, body) in enumerate(steps):
        x = Inches(0.15 + i * 2.63)
        rect(sl, x, Inches(1.45), Inches(2.45), Inches(3.2), color)
        txt(sl, title, x + Inches(0.1), Inches(1.5), Inches(2.25), Inches(0.7),
            font_size=13, bold=True, color=C_WHITE)
        txt(sl, body, x + Inches(0.1), Inches(2.2), Inches(2.25), Inches(2.3),
            font_size=10.5, color=RGBColor(0xcc, 0xdd, 0xff))
        if i < 4:
            txt(sl, "→", x + Inches(2.45), Inches(2.8), Inches(0.2), Inches(0.4),
                font_size=18, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)

    # Bottom: metrics
    rect(sl, Inches(0.3), Inches(4.85), Inches(12.73), Inches(2.35), C_WHITE)
    rect(sl, Inches(0.3), Inches(4.85), Inches(12.73), Inches(0.4), C_DARK_BLUE)
    txt(sl, "Метрики качества ML-моделей", Inches(0.5), Inches(4.9), Inches(6), Inches(0.3),
        font_size=12, bold=True, color=C_WHITE)

    metrics = [
        ("Prophet Forecast", "MAE / RMSE / MAPE", "Backtesting: fitted vs actual.\nНа мок-данных — in-sample. В продакшене — walk-forward CV"),
        ("KMeans Кластеры", "Silhouette = 0.835", "От -1 до 1. >0.5 = хорошо, >0.7 = сильные кластеры.\nDavies-Bouldin = 1.204 (чем ниже тем лучше)"),
        ("Корреляция", "Pearson r = 0.415", "Нарушения Коргау ↔ инциденты. Умеренная связь (17% дисперсии).\np-value < 0.05 — статистически значимо"),
        ("Risk Scoring", "0–100 (4 уровня)", "Формула ТЗ: 0.35×incidents + 0.30×violations\n+ 0.20×trend + 0.15×severity. Min-max нормализация"),
    ]
    for i, (name, value, desc) in enumerate(metrics):
        lx = Inches(0.4 + (i % 2) * 6.4)
        ty = Inches(5.35 + (i // 2) * 0.9)
        txt(sl, name, lx, ty, Inches(2.0), Inches(0.3), font_size=11, bold=True, color=C_DARK_BLUE)
        txt(sl, value, lx + Inches(2.0), ty, Inches(1.8), Inches(0.3), font_size=11, bold=True, color=C_GREEN)
        txt(sl, desc, lx + Inches(3.8), ty, Inches(2.5), Inches(0.8), font_size=9, color=C_GRAY)
    footer(sl, 5, TOTAL)


def slide_why_methods(prs):
    """Почему выбраны именно эти методы."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "Почему эти методы? Альтернативы и трейдоффы",
               "Обоснование выбора Prophet, KMeans, Gemini, Risk Scoring")

    comparisons = [
        (C_DARK_BLUE, "Прогноз: Prophet",
         "Выбрали Prophet\n• Работает с <36 месяцами данных\n• Автоматические changepoints\n• Встроенный 95% CI\n• Устойчив к пропускам",
         "Альтернативы:\n• ARIMA — требует стационарность\n• LSTM — нужно >100 точек\n• Exp. Smoothing — нет CI\n• XGBoost — нет временной оси"),
        (RGBColor(0x6a, 0x1b, 0x9a), "Кластеры: TF-IDF + KMeans",
         "Выбрали TF-IDF + KMeans\n• Интерпретируемо: top-terms\n• Быстро: <1 сек на 500 текстов\n• Silhouette 0.835 (отлично)\n• Не нужна разметка",
         "Альтернативы:\n• BERT embeddings — тяжёлый\n• LDA (topic model) — хуже для\n  коротких текстов\n• DBSCAN — не контролируемое K"),
        (C_RED, "NLP: Gemini Flash",
         "Выбрали Gemini LLM\n• Zero-shot: без обучения\n• Понимает русский + домен ОТ\n• Классификация + рекомендации\n  + CV-анализ фото в одной модели",
         "Альтернативы:\n• Fine-tuned BERT — нужна разметка\n  (~1000 примеров на класс)\n• spaCy NER — только сущности\n• Regex — не масштабируется"),
        (C_GREEN, "Сценарий: Pearson × меры",
         "Выбрали корреляцию + формулу\n• Прозрачно для бизнеса\n• Pearson r=0.415 (p<0.05)\n• 9 мер × вес категории\n• Экономика: × 5 млн ₸/НС",
         "Ограничения:\n• Корреляция ≠ причинность\n• r=0.415 = 17% дисперсии\n• Веса мер — экспертные\n• В продакшене → causal inference"),
    ]

    for i, (color, title, chosen, alt) in enumerate(comparisons):
        col, row = i % 2, i // 2
        lx = Inches(0.3 + col * 6.52)
        ty = Inches(1.45 + row * 2.9)

        rect(sl, lx, ty, Inches(6.22), Inches(2.7), C_WHITE)
        rect(sl, lx, ty, Inches(6.22), Inches(0.45), color)
        txt(sl, title, lx + Inches(0.15), ty + Inches(0.07), Inches(5.8), Inches(0.32),
            font_size=13, bold=True, color=C_WHITE)

        txt(sl, chosen, lx + Inches(0.15), ty + Inches(0.55), Inches(3.0), Inches(2.0),
            font_size=10, color=C_DARK)
        txt(sl, alt, lx + Inches(3.2), ty + Inches(0.55), Inches(2.9), Inches(2.0),
            font_size=10, color=C_GRAY, italic=True)


def slide_production(prs):
    """Что меняется для продакшн данных."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "Мок-данные → Продакшн: что меняется",
               "Хакатон = proof of concept. Продакшн = другой уровень валидации")

    items = [
        (C_AMBER, "Сейчас (мок)", "Продакшн",
         "220 инцидентов, <36 месяцев",
         "1000+ инцидентов, 5+ лет истории → Prophet с годовой + недельной сезонностью"),
        (C_AMBER, "Сейчас (мок)", "Продакшн",
         "In-sample backtesting (MAE на train)",
         "Walk-forward CV: train 80% → predict 20% → реальный MAE/RMSE/MAPE"),
        (C_AMBER, "Сейчас (мок)", "Продакшн",
         "k=6 кластеров (хардкод)",
         "Elbow method + Silhouette sweep (k=2..15) → автовыбор оптимального K"),
        (C_AMBER, "Сейчас (мок)", "Продакшн",
         "Risk Score: экспертные веса 0.35/0.30/0.20/0.15",
         "Logistic Regression / XGBoost обучен на экспертной разметке → SHAP feature importance"),
        (C_AMBER, "Сейчас (мок)", "Продакшн",
         "Gemini без валидации (zero-shot)",
         "Тест на 100+ размеченных примерах: Precision/Recall/F1 > 0.85 для каждого типа"),
        (C_AMBER, "Сейчас (мок)", "Продакшн",
         "Pearson r=0.415 как causal",
         "Partial correlation (контроль на размер орг) + Granger causality test + time-lag анализ"),
    ]

    for i, (color, tag_now, tag_prod, now_text, prod_text) in enumerate(items):
        ty = Inches(1.45 + i * 0.95)
        rect(sl, Inches(0.3), ty, Inches(12.73), Inches(0.85), C_WHITE if i % 2 == 0 else RGBColor(0xf5, 0xf8, 0xff))

        rect(sl, Inches(0.3), ty, Inches(0.08), Inches(0.85), color)
        txt(sl, now_text, Inches(0.5), ty + Inches(0.08), Inches(5.5), Inches(0.65),
            font_size=11, color=C_DARK)
        txt(sl, "→", Inches(6.1), ty + Inches(0.2), Inches(0.3), Inches(0.4),
            font_size=16, bold=True, color=C_GREEN)
        txt(sl, prod_text, Inches(6.5), ty + Inches(0.08), Inches(6.3), Inches(0.65),
            font_size=11, color=C_GREEN, bold=True)
    footer(sl, 9, TOTAL)


def slide_glossary(prs):
    """Глоссарий ML-терминов простым языком."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_LIGHT)
    header_bar(sl, "ML/AI — что это значит простым языком",
               "Объяснение ключевых терминов для не-технических слушателей")

    terms = [
        (C_DARK_BLUE, "Prophet",
         "Модель Facebook для прогноза временных рядов. Как Excel-тренд, но умнее —\n"
         "учитывает сезонность (зима/лето), тренд (растёт/падает) и аномалии."),
        (C_BLUE, "TF-IDF + KMeans",
         "TF-IDF: превращает текст в числа — считает «насколько важно это слово в документе».\n"
         "KMeans: группирует похожие тексты в кластеры. Как сортировка писем по темам."),
        (C_GREEN, "Pearson r (корреляция)",
         "Число от -1 до +1. Показывает связь двух величин. r=0.4 значит «когда нарушений\n"
         "больше — инцидентов тоже больше, но не всегда». Это связь, не причина!"),
        (C_AMBER, "Silhouette Score",
         "Оценка качества кластеризации от -1 до 1. Наш 0.835 = «кластеры чёткие, группы\n"
         "хорошо разделены». >0.7 = отлично, <0.3 = плохо."),
        (C_RED, "MAE / RMSE / MAPE",
         "Метрики ошибки прогноза. MAE: «в среднем ошибаемся на X инцидентов».\n"
         "RMSE: то же, но штрафует большие ошибки сильнее. MAPE: ошибка в процентах."),
        (RGBColor(0x6a, 0x1b, 0x9a), "Confidence Interval (95% CI)",
         "Коридор прогноза: «мы на 95% уверены что реальное значение будет в этих границах».\n"
         "Узкий CI = уверенный прогноз. Широкий CI = большая неопределённость."),
        (C_DARK_BLUE, "Zero-shot (Gemini)",
         "AI классифицирует текст БЕЗ обучения на наших данных. Мы просто описываем задачу\n"
         "в промпте, и LLM применяет свои знания. Быстро, но нужно валидировать."),
        (C_BLUE, "Risk Scoring",
         "Формула: каждой организации присваивается балл 0–100 на основе 4 факторов:\n"
         "частота инцидентов (35%) + нарушения (30%) + тренд (20%) + тяжесть (15%)."),
    ]

    for i, (color, term, desc) in enumerate(terms):
        col, row = i % 2, i // 2
        lx = Inches(0.3 + col * 6.52)
        ty = Inches(1.45 + row * 1.45)

        rect(sl, lx, ty, Inches(6.22), Inches(1.3), C_WHITE)
        rect(sl, lx, ty, Inches(0.08), Inches(1.3), color)
        txt(sl, term, lx + Inches(0.2), ty + Inches(0.05), Inches(2.0), Inches(0.3),
            font_size=13, bold=True, color=color)
        txt(sl, desc, lx + Inches(0.2), ty + Inches(0.38), Inches(5.8), Inches(0.85),
            font_size=10, color=C_DARK)
    footer(sl, 4, TOTAL)


TOTAL = 10


def slide_qr(prs):
    """QR код для live демо."""
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, C_DARK_BLUE)
    rect(sl, 0, Inches(6.2), W, Inches(0.15), C_AMBER)

    txt(sl, "Попробуйте прямо сейчас", Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7),
        font_size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "Откройте камеру телефона и наведите на QR-код",
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.4),
        font_size=15, color=C_AMBER, align=PP_ALIGN.CENTER)

    try:
        sl.shapes.add_picture("qr_streamlit.png",
                              Inches(4.67), Inches(1.7), width=Inches(4.0), height=Inches(4.0))
    except Exception:
        rect(sl, Inches(4.67), Inches(1.7), Inches(4.0), Inches(4.0), C_WHITE)

    txt(sl, "http://89.207.255.254:8502",
        Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.4),
        font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "Инциденты · Карты Коргау · Предикт · Алерты · PDF/Excel",
        Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.35),
        font_size=12, color=RGBColor(0xaa, 0xcc, 0xff), align=PP_ALIGN.CENTER)
    footer(sl, 9, TOTAL)


def build():
    prs = new_prs()

    # 3-min pitch: 10 слайдов
    # 0:00-0:30  Титул + Проблема
    slide_title(prs)             # 1
    slide_problem(prs)           # 2

    # 0:30-1:00  Решение + ML
    slide_solution(prs)          # 3
    slide_glossary(prs)          # 4 — ML термины простым языком
    slide_ml_pipeline(prs)       # 5 — pipeline + метрики

    # 1:00-2:00  LIVE ДЕМО (переключиться на Streamlit)
    slide_qr(prs)                # 6 — QR код → жюри открывает на телефоне

    # 2:00-2:30  Экономика + критерии
    slide_economic(prs)          # 7
    slide_criteria(prs)          # 8

    # 2:30-3:00  Продакшн + Спасибо
    slide_production(prs)        # 9
    slide_results(prs)           # 10

    out = "HSE_Presentation.pptx"
    prs.save(out)
    print(f"✅  Презентация сохранена: {out}  ({len(prs.slides)} слайдов)")


if __name__ == "__main__":
    build()